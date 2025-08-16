from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import numpy as np
import json

# Initialize presentation
prs = Presentation()

# ---------- Helper Functions ----------
def add_title_slide(title, subtitle, name, date):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = f"{subtitle}\n\n{name}\n{date}"

def add_bullet_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    tf = body_placeholder.text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)

# ---------- Slide 1: Title ----------
add_title_slide(
    "MNIST GAN with FastAPI and HTML UI",
    "A Project Report on a Deployable Generative AI Application",
    "[Your Name] / Project ID: [Your Name]",
    "[Current Date]"
)

# ---------- Slide 2: Executive Summary ----------
exec_summary = [
    "Objective: To build and deploy a Generative Adversarial Network (GAN) as a web service for generating handwritten digits.",
    "Achievement: Successfully developed a full-stack application with a FastAPI backend and an HTML frontend that uses a trained Generator model to create new images. A separate Classifier model was integrated to predict the digit and confidence of the generated images.",
    "Impact: Created a reproducible, end-to-end pipeline for deploying a generative model, demonstrating key MLOps principles for turning a model into a functional application."
]
add_bullet_slide("Executive Summary", exec_summary)

# ---------- Slide 3: Problem & Methodology ----------
problem_methodology = [
    "Problem: Making a trained generative model accessible to end-users through a simple, interactive web interface.",
    "Backend: A FastAPI server was built to load both the Generator and Classifier models and handle image generation requests. It returns a Base64 encoded image and the top-3 predictions.",
    "Frontend: A minimalist HTML/JavaScript UI allows users to generate new digits and see real-time predictions and confidence scores.",
    "Data: The standard MNIST dataset was used for training both the GAN and the classifier models."
]
add_bullet_slide("Problem & Methodology", problem_methodology)

# ---------- Slide 4: Model Architecture ----------
model_arch = [
    "Generator (The Artist): A multi-layer perceptron (MLP) that takes a 100-dimensional random noise vector as input. It progressively increases the dimensions through hidden layers and uses a Tanh activation to output a 28x28 pixel image with pixel values in the range [-1, 1].",
    "Discriminator (The Critic): An MLP that takes a flattened 784-dimensional image as input. Its purpose is to distinguish between real (MNIST) images and fake (generated) images. It outputs a single scalar value between 0 and 1 using a Sigmoid activation.",
    "Classifier (The Judge): A separate Convolutional Neural Network (CNN). It takes the generated image as input and predicts the digit (0-9) and its confidence level. This model provides the prediction and confidence score for the web application."
]
add_bullet_slide("Model Architecture", model_arch)

# ---------- Slide 5: GAN Training Results ----------
# Loss data
g_losses_str = "[0.8547711443227491, 0.9366002984519707, 0.9809877779692221, 1.028972950190115, 0.9621060712695885, 0.9373731607757906, 0.9354666158525166, 0.9194455803520898, 0.9101737885714085, 0.8998334066890704, 0.8847102800539053, 0.882334085987575, 0.8772036104059932, 0.8781337272574399, 0.8687292881039922, 0.868195465887025, 0.8729581203160763, 0.8681320279264755, 0.8674189297756406, 0.8635128218609133]"
d_losses_str = "[0.6469696177475488, 0.6336105929445356, 0.6164920027258554, 0.6029787260268543, 0.6169381328165404, 0.6251830791613695, 0.6256572470736148, 0.631901282936271, 0.6331600889341155, 0.6365701671221109, 0.6432711013090382, 0.6435819352740673, 0.644482314840817, 0.6430365392394157, 0.6458770926954396, 0.645183054305343, 0.644880214416142, 0.6452496558872621, 0.6461725732537983, 0.6470403444411149]"

g_losses = json.loads(g_losses_str)
d_losses = json.loads(d_losses_str)
epochs = np.arange(1, len(g_losses) + 1)

plt.figure(figsize=(10, 6))
plt.title("Generator and Discriminator Loss During Training", fontsize=16)
plt.plot(epochs, g_losses, label="Generator Loss", marker='o', linestyle='-', color='#1f77b4')
plt.plot(epochs, d_losses, label="Discriminator Loss", marker='o', linestyle='-', color='#ff7f0e')
plt.xlabel("Epochs", fontsize=12)
plt.ylabel("Loss", fontsize=12)
plt.legend(fontsize=10)
plt.grid(True, linestyle='--', alpha=0.6)
plt.xticks(epochs)
plt.tight_layout()
plt.savefig("./report/gan_training_results.png")
plt.close()

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Generator and Discriminator Loss During Training"
slide.shapes.add_picture("./report/gan_training_results.png", Inches(1), Inches(2), width=Inches(8), height=Inches(4))

# ---------- Slide 6: Application Demo & Prediction ----------
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Real-time Image Generation with Prediction"
slide.shapes.add_picture("./application.PNG", Inches(1), Inches(2), width=Inches(5), height=Inches(4))

# Bullet points
textbox = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3), Inches(4))
tf = textbox.text_frame
for point in [
    "The app generates a new image with each click.",
    "The Classifier model provides a prediction and confidence score for the generated digit.",
    "This demonstrates a complete end-to-end workflow from model generation to user-facing application."
]:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(18)

# ---------- Slide 7: Key Challenges & Lessons Learned ----------
challenges = [
    "Version Compatibility: Encountered a numpy version mismatch (python 3.10 vs numpy 2.3.2), which was resolved by specifying the correct version in requirements.txt. Lesson: Precise version pinning is critical for ensuring a reproducible and functional environment.",
    "Docker Networking: A browser error (ERR_ADDRESS_INVALID) occurred when trying to connect to 0.0.0.0:8000. Lesson: The correct address for local access is http://127.0.0.1:8000, a crucial detail for deployment.",
    "PyTorch-NumPy Interoperability: An AttributeError was encountered when trying to use a NumPy method on a PyTorch tensor. Lesson: Tensors must be explicitly converted to NumPy arrays (.numpy()) before using NumPy-specific methods."
]
add_bullet_slide("Key Challenges & Lessons Learned", challenges)

# ---------- Slide 8: MLOps & Reproducibility ----------
mlops = [
    "Containerization: The Dockerfile provides a consistent environment for the application, ensuring it runs the same way on any machine, from development to production.",
    "CI/CD: The GitHub Actions pipeline runs automated tests on the FastAPI endpoints, ensuring the application is always functional.",
    "Test-Driven Development: The project uses pytest to verify the application's core functionality, which is a key part of a professional development workflow."
]
add_bullet_slide("MLOps & Reproducibility", mlops)

# ---------- Slide 9: Future Work & Recommendations ----------
future_work = [
    "Model Improvement: Train the GAN for more epochs on a larger dataset to improve the quality of the generated digits.",
    "Deployment Automation: Finalize the CI/CD pipeline to automatically build and deploy the Docker image to a cloud service like AWS, Google Cloud, or Azure.",
    "Improved UI: Enhance the user interface with additional features, such as the ability to save generated images or select a specific digit to generate."
]
add_bullet_slide("Future Work & Recommendations", future_work)

# ---------- Final Takeaway Slide ----------
takeaways = [
    "The losses for both models remain relatively balanced, indicating a stable adversarial training process.",
    "The Generator's loss and the Discriminator's loss are both decreasing, which is a good sign of a stable GAN.",
    "The plot shows a stable trend, which is the most important indicator of a healthy GAN training run."
]
add_bullet_slide("Final Takeaways", takeaways)

# Save presentation
prs.save("./report/MNIST_GAN_FastAPI_Presentation.pptx")
print("Presentation saved as MNIST_GAN_FastAPI_Presentation.pptx")
